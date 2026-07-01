from pathlib import Path
from flask import Flask, jsonify, request, send_from_directory, send_file
from flask_cors import CORS
import json
import os
from pptx import Presentation
from datetime import datetime
import io

BASE_DIR = Path(__file__).resolve().parent
FROND_DIR = (BASE_DIR.parent / 'Frond_end').resolve()
STORE_FILE = BASE_DIR / 'data_store.json'

PERFILES = [
    {
        'id': 1,
        'nombre': 'ADMINISTRADOR BASE DATOS ESTANDAR',
        'costoEmpresa': 9208000,
        'uuid': '258937be-b770-4123-a9df-3f757f513919'
    },
    {
        'id': 2,
        'nombre': 'ADMINISTRADOR BASE DATOS JUNIOR',
        'costoEmpresa': 3400000,
        'uuid': '65fa0b79-b6e1-4874-bcbe-4492b866f92b'
    },
    {
        'id': 3,
        'nombre': 'ADMINISTRADOR BASE DATOS L3',
        'costoEmpresa': 16900000,
        'uuid': '01454e77-4e48-43c0-9697-245caeb5fd4f'
    },
    {
        'id': 4,
        'nombre': 'ADMINISTRADOR BASE DE DATOS N2',
        'costoEmpresa': 9842600,
        'uuid': '5b318d1c-724b-4594-b62c-74885f787bb0'
    },
    {
        'id': 5,
        'nombre': 'ADMINISTRADOR BASES DATOS JR',
        'costoEmpresa': 3097500,
        'uuid': '83a6d230-3229-4587-a3e2-cca0d5df7a34'
    },
    {
        'id': 6,
        'nombre': 'ADMINISTRADOR BASES DATOS JUNIOR',
        'costoEmpresa': 3038000,
        'uuid': 'aa3b6575-ca13-4639-879f-953b876b978c'
    },
    {
        'id': 7,
        'nombre': 'ADMINISTRADOR BASES DE DATOS',
        'costoEmpresa': 6157250,
        'uuid': '9f80d540-da0c-47c4-991f-787756b0fd86'
    },
    {
        'id': 8,
        'nombre': 'ADMINISTRADOR BASES DE DATOS ESTANDAR',
        'costoEmpresa': 8050000,
        'uuid': 'e18838c6-dcb5-4e58-92dc-1aad8f8461ee'
    },
    {
        'id': 9,
        'nombre': 'ADMINISTRADOR BASES DE DATOS JUNIOR',
        'costoEmpresa': 5245333,
        'uuid': 'd74b53f6-55e4-4710-9a94-942bb6742fac'
    },
    {
        'id': 10,
        'nombre': 'ADMINISTRADOR BASES DE DATOS N1',
        'costoEmpresa': 6100000,
        'uuid': '379bd29d-0719-4ef0-8d05-60db30159455'
    },
    {
        'id': 11,
        'nombre': 'ADMINISTRADOR BASES DE DATOS SENIOR',
        'costoEmpresa': 8050000,
        'uuid': '29ea996f-8052-4bca-a26b-baa8550a2731'
    },
    {
        'id': 12,
        'nombre': 'ANALISTA BI',
        'costoEmpresa': 4648000,
        'uuid': '6b384b90-0ebf-490e-a49e-29d5796deb3a'
    },
    {
        'id': 13,
        'nombre': 'ANALISTA BI ESTANDAR',
        'costoEmpresa': 4147000,
        'uuid': '89b19a7d-17bd-46f0-a967-43725bd2ccf9'
    },
    {
        'id': 14,
        'nombre': 'ANALISTA AUTOMATIZACION ESTANDAR',
        'costoEmpresa': 5533750,
        'uuid': '4d77a2f9-19f2-4448-b2f1-c69f5bb82d91'
    },
    {
        'id': 15,
        'nombre': 'ANALISTA AUTOMATIZACION JUNIOR',
        'costoEmpresa': 4091333,
        'uuid': '28ff8b64-175b-403e-b167-7f711487c703'
    },
    {
        'id': 16,
        'nombre': 'ANALISTA AUTOMATIZACION SENIOR',
        'costoEmpresa': 6922000,
        'uuid': 'b41654e8-b395-4abb-a13b-16675ede7465'
    },
    {
        'id': 17,
        'nombre': 'ANALISTA ASEGURAMIENTO CALIDAD',
        'costoEmpresa': 2662000,
        'uuid': 'f004a1c8-c87e-4956-b56e-265cc0a177f4'
    },
    {
        'id': 18,
        'nombre': 'ANALISTA ABAP',
        'costoEmpresa': 3813000,
        'uuid': '71b00239-9970-4769-8d65-d8bc0e4212c8'
    },
    {
        'id': 19,
        'nombre': 'ADMINISTRADOR PROYECTOS SENIOR BILINGÜE',
        'costoEmpresa': 12000000,
        'uuid': 'e5db0d17-8741-4e4a-b79a-bac07a412c9a'
    }
]

app = Flask(__name__, static_folder=str(FROND_DIR), static_url_path='')
CORS(app)


def load_store():
    if not STORE_FILE.exists():
        return {'clienteDatos': {}, 'perfilesSeleccionadosFinal': [], 'perfilesDescripcion': []}
    with STORE_FILE.open('r', encoding='utf-8') as f:
        return json.load(f)


def save_store(data):
    with STORE_FILE.open('w', encoding='utf-8') as f:
        json.dump(data, f, indent=2, ensure_ascii=False)


@app.route('/')
def index():
    return app.send_static_file('Login.html')


@app.route('/<path:path>')
def static_files(path):
    return send_from_directory(str(FROND_DIR), path)


@app.route('/api/perfiles', methods=['GET'])
def get_perfiles():
    q = request.args.get('q', '').strip().lower()
    if q:
        filtered = [p for p in PERFILES if q in p['nombre'].lower()]
    else:
        filtered = PERFILES
    return jsonify(filtered)


@app.route('/api/cliente', methods=['POST'])
def save_cliente():
    payload = request.get_json(silent=True) or {}
    store = load_store()
    store['clienteDatos'] = payload
    save_store(store)
    return jsonify({'ok': True, 'clienteDatos': payload})


@app.route('/api/seleccionados', methods=['POST'])
def save_seleccionados():
    payload = request.get_json(silent=True) or {}
    seleccionados = payload.get('perfiles', [])
    store = load_store()
    store['perfilesSeleccionadosFinal'] = seleccionados
    save_store(store)
    return jsonify({'ok': True, 'perfilesSeleccionadosFinal': seleccionados})


@app.route('/api/descripcion', methods=['POST'])
def save_descripcion():
    payload = request.get_json(silent=True) or {}
    descripcion = payload.get('perfiles', [])
    store = load_store()
    store['perfilesDescripcion'] = descripcion
    save_store(store)
    return jsonify({'ok': True, 'perfilesDescripcion': descripcion})


@app.route('/api/reporte', methods=['GET'])
def get_reporte():
    store = load_store()
    cliente = store.get('clienteDatos', {})
    perfiles = store.get('perfilesSeleccionadosFinal', [])
    descripcion = store.get('perfilesDescripcion', [])

    perfiles_map = {p.get('uuid') or p.get('nombre'): p for p in descripcion}
    reporte_perfiles = []

    for perfil in perfiles:
        key = perfil.get('uuid') or perfil.get('nombre')
        descr = perfiles_map.get(key, {})
        reporte_perfiles.append({
            'uuid': perfil.get('uuid'),
            'nombre': perfil.get('nombre'),
            'costoEmpresa': perfil.get('costoEmpresa'),
            'margen': perfil.get('margen'),
            'costoMas': perfil.get('costoMas'),
            'experiencia': descr.get('experiencia', ''),
            'descripcion': descr.get('descripcion', ''),
            'conocimientos': descr.get('conocimientos', ''),
            'funciones': descr.get('funciones', '')
        })

    return jsonify({
        'clienteDatos': cliente,
        'perfiles': reporte_perfiles
    })


@app.route('/api/reporte/pptx', methods=['GET'])
def export_pptx():
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    store = load_store()
    cliente = store.get('clienteDatos', {})
    perfiles = store.get('perfilesSeleccionadosFinal', [])
    
    try:
        # Crear presentación nueva
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Portada
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(45, 90, 168)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = "OFERTA COMERCIAL"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        info_frame = info_box.text_frame
        info_frame.word_wrap = True
        p = info_frame.paragraphs[0]
        p.text = f"Cliente: {cliente.get('nombreCliente', 'N/A')}\nFecha: {datetime.now().strftime('%d-%m-%Y')}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Slide 2: Datos del Cliente
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Datos del Cliente"
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        data = [
            f"Cliente: {cliente.get('nombreCliente', 'N/A')}",
            f"NIT: {cliente.get('nit', 'N/A')}",
            f"Correo: {cliente.get('correocorporativo', 'N/A')}",
            f"Representante: {cliente.get('representante', 'N/A')}",
            f"Comercial: {cliente.get('nombreComercial', 'N/A')}",
            f"Preventa: {cliente.get('nombrePreventa', 'N/A')}",
            f"Tipo Contrato: {cliente.get('tipoContrato', 'N/A')}",
            f"Meses: {cliente.get('mesesContrato', 'N/A')}"
        ]
        
        for i, text in enumerate(data):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(14)
            p.level = 0
        
        # Slide 3: Perfiles
        if perfiles:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = "Perfiles Seleccionados"
            
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5.5)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            
            for perfil in perfiles:
                p = tf.add_paragraph()
                p.text = f"• {perfil.get('nombre', 'N/A')} - COP {perfil.get('costoEmpresa', 0):,}"
                p.font.size = Pt(12)
                p.level = 0
        
        # Guardar en BytesIO
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        
        from flask import Response
        
        response = Response(pptx_io.getvalue())
        response.headers['Content-Disposition'] = 'attachment; filename=Oferta.pptx'
        response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        response.headers['Content-Length'] = len(pptx_io.getvalue())
        
        return response
    except Exception as e:
        print(f"Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': str(e)}), 500


if __name__ == '__main__':
    if not STORE_FILE.exists():
        save_store({'clienteDatos': {}, 'perfilesSeleccionadosFinal': [], 'perfilesDescripcion': []})
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=False)
